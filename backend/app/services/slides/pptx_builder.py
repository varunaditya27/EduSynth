"""
PPTX builder with visual themes and diagram support.
(Aesthetic-max version with gradients, accents, and simple diagrams.)
"""

from io import BytesIO
from typing import Dict, Any, Tuple, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

from app.schemas.slides import LecturePlan, SlideItem
from .diagram_draw import hex_to_rgb
from .icons import get_point_icon



# ---------- Color helpers (python-pptx solid fills don't support alpha) ----------

def _blend_with_bg(
    rgb: tuple[int, int, int],
    bg: tuple[int, int, int] = (255, 255, 255),
    alpha: float = 0.5,
) -> tuple[int, int, int]:
    """
    Approximate opacity by blending 'rgb' over 'bg'.
    alpha=0.0 -> fully bg; alpha=1.0 -> original rgb.
    """
    a = max(0.0, min(1.0, float(alpha)))
    return tuple(int(round(a * rgb[i] + (1.0 - a) * bg[i])) for i in range(3))


# ---------- Gradient background renderer (Pillow) ----------

def _render_gradient_bg(
    width_px: int,
    height_px: int,
    stops: Tuple[str, str],
    mode: str = "linear",  # "linear" | "radial"
) -> bytes:
    """
    Render a simple gradient background as PNG bytes.
    """
    img = Image.new("RGB", (width_px, height_px))
    draw = ImageDraw.Draw(img)

    color1 = hex_to_rgb(stops[0])
    color2 = hex_to_rgb(stops[1])

    if mode == "radial":
        # Radial from center
        cx, cy = width_px // 2, height_px // 2
        max_r = ((width_px ** 2 + height_px ** 2) ** 0.5) / 2
        for y in range(height_px):
            for x in range(width_px):
                dx, dy = x - cx, y - cy
                dist = (dx * dx + dy * dy) ** 0.5
                t = min(dist / max_r, 1.0)
                r = int(color1[0] * (1 - t) + color2[0] * t)
                g = int(color1[1] * (1 - t) + color2[1] * t)
                b = int(color1[2] * (1 - t) + color2[2] * t)
                draw.point((x, y), fill=(r, g, b))
    else:
        # Linear top->bottom
        for y in range(height_px):
            t = y / max(1, height_px - 1)
            r = int(color1[0] * (1 - t) + color2[0] * t)
            g = int(color1[1] * (1 - t) + color2[1] * t)
            b = int(color1[2] * (1 - t) + color2[2] * t)
            draw.line([(0, y), (width_px, y)], fill=(r, g, b))

    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------- (Optional) shadow helper (creation-order "shadow") ----------

def _add_shadow_shape(shapes, shape, theme: Dict[str, Any]) -> None:
    """
    Fake a shadow by duplicating the shape behind with a slight offset.
    """
    offset_pt = theme["depth"]["shape_shadow_offset_pt"]
    # opacity not directly supported; we just use black behind
    shadow = shapes.add_shape(
        shape.shape_type,
        shape.left + Pt(offset_pt),
        shape.top + Pt(offset_pt),
        shape.width,
        shape.height,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shadow.line.color.rgb = RGBColor(0, 0, 0)
    # No z-order API; rely on creation order (shadow created after goes on top),
    # so if you need it truly behind, create shadow first, then the main shape.


# ---------- Layout selector ----------

def pick_layout(slide: SlideItem) -> str:
    """
    Return the layout type for a slide.
    """
    if slide.diagram:
        return slide.diagram
    return "text"


# ---------- Diagram renderers ----------

def add_process_diagram(
    shapes,
    left: float,
    top: float,
    width: float,
    height: float,
    points: List[str],
    theme: Dict[str, Any],
) -> None:
    """
    Left-to-right rounded boxes + connectors for "process".
    """
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    steps = max(1, min(len(points), 5))
    box_w = width / (steps * 1.5)
    spacing = (width - steps * box_w) / (steps - 1) if steps > 1 else 0.0

    accent_rgb = hex_to_rgb(theme["colors"]["accent"])
    text_rgb = (255, 255, 255)

    for i, label in enumerate(points[:steps]):
        box_left = left + i * (box_w + spacing)
        box_top = top + (height - Inches(1.0)) / 2

        rect = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box_left),
            Inches(box_top),
            Inches(box_w),
            Inches(1.0),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(*accent_rgb)
        rect.line.color.rgb = RGBColor(*accent_rgb)

        tf = rect.text_frame
        tf.text = label
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)

        # connector arrow to next
        if i < steps - 1:
            x1 = box_left + box_w
            y = box_top + 0.5
            conn = shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x1),
                Inches(y),
                Inches(x1 + spacing),
                Inches(y),
            )
            conn.line.color.rgb = RGBColor(*accent_rgb)
            conn.line.width = Pt(2)


def add_tree_diagram(
    shapes,
    left: float,
    top: float,
    width: float,
    height: float,
    points: List[str],
    theme: Dict[str, Any],
) -> None:
    """
    Parent box with up to 3 children under it for "tree".
    """
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    accent = hex_to_rgb(theme["colors"]["accent"])
    accent2 = hex_to_rgb(theme["colors"]["accent2"])
    text_white = (255, 255, 255)

    # Parent
    parent_w = width * 0.4
    parent_l = left + (width - parent_w) / 2
    parent = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(parent_l),
        Inches(top),
        Inches(parent_w),
        Inches(0.8),
    )
    parent.fill.solid()
    parent.fill.fore_color.rgb = RGBColor(*accent)
    parent.line.color.rgb = RGBColor(*accent)

    tf = parent.text_frame
    tf.text = points[0] if points else "Parent"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*text_white)

    # Children (max 3)
    children = points[1:4]
    if children:
        child_w = width / (len(children) + 1)
        child_top = top + 2.0
        for i, label in enumerate(children):
            child_left = left + i * child_w + child_w * 0.25
            child = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(child_left),
                Inches(child_top),
                Inches(child_w * 0.5),
                Inches(0.6),
            )
            child.fill.solid()
            child.fill.fore_color.rgb = RGBColor(*accent2)
            child.line.color.rgb = RGBColor(*accent2)
            tf = child.text_frame
            tf.text = label
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(*text_white)

            # connector from parent bottom center to child top center
            conn = shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                parent.left + parent.width / 2,
                parent.top + parent.height,
                child.left + child.width / 2,
                child.top,
            )
            conn.line.color.rgb = RGBColor(*accent2)
            conn.line.width = Pt(2)


def add_timeline_diagram(
    shapes,
    left: float,
    top: float,
    width: float,
    height: float,
    points: List[str],
    theme: Dict[str, Any],
) -> None:
    """
    Horizontal timeline with circular markers.
    """
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    accent = hex_to_rgb(theme["colors"]["accent"])
    text_rgb = hex_to_rgb(theme["colors"]["text"])

    y = top + height / 2
    line = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left),
        Inches(y),
        Inches(left + width),
        Inches(y),
    )
    line.line.color.rgb = RGBColor(*accent)
    line.line.width = Pt(3)

    n = max(1, min(len(points), 5))
    gap = width / (n + 1)
    for i, label in enumerate(points[:n]):
        cx = left + (i + 1) * gap
        # marker circle
        marker = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - 0.15),
            Inches(y - 0.15),
            Inches(0.3),
            Inches(0.3),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(*accent)
        marker.line.color.rgb = RGBColor(*accent)

        # label
        tb = shapes.add_textbox(
            Inches(cx - 0.6),
            Inches(y + 0.35),
            Inches(1.2),
            Inches(0.6),
        )
        tf = tb.text_frame
        tf.text = label
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*text_rgb)


def add_compare_layout(
    shapes,
    left: float,
    top: float,
    width: float,
    height: float,
    points: List[str],
    theme: Dict[str, Any],
) -> None:
    """
    Two side-by-side panels for "compare".
    """
    from pptx.enum.shapes import MSO_SHAPE

    accent = hex_to_rgb(theme["colors"]["accent"])
    accent2 = hex_to_rgb(theme["colors"]["accent2"])
    text_rgb = hex_to_rgb(theme["colors"]["text"])

    panel_w = (width - 0.2) / 2

    # Split points
    mid = len(points) // 2
    left_points = points[:mid] if mid > 0 else points[:2]
    right_points = points[mid:] if mid > 0 else points[2:]

    # Left panel
    left_panel = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(panel_w),
        Inches(height),
    )
    left_panel.fill.solid()
    # Fake 50% opacity over white
    soft_accent = _blend_with_bg(accent, (255, 255, 255), alpha=0.5)
    left_panel.fill.fore_color.rgb = RGBColor(*soft_accent)
    left_panel.line.color.rgb = RGBColor(*accent)

    # Right panel
    right_panel = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + panel_w + 0.2),
        Inches(top),
        Inches(panel_w),
        Inches(height),
    )
    right_panel.fill.solid()
    # Fake 40% opacity over white
    soft_accent2 = _blend_with_bg(accent2, (255, 255, 255), alpha=0.4)
    right_panel.fill.fore_color.rgb = RGBColor(*soft_accent2)
    right_panel.line.color.rgb = RGBColor(*accent2)

    # Text blocks
    for i, point in enumerate(left_points[:3]):
        tb = shapes.add_textbox(
            Inches(left + 0.2),
            Inches(top + 0.25 + i * 0.65),
            Inches(panel_w - 0.4),
            Inches(0.55),
        )
        tf = tb.text_frame
        tf.text = f"• {point}"
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)

    for i, point in enumerate(right_points[:3]):
        tb = shapes.add_textbox(
            Inches(left + panel_w + 0.4),
            Inches(top + 0.25 + i * 0.65),
            Inches(panel_w - 0.4),
            Inches(0.55),
        )
        tf = tb.text_frame
        tf.text = f"• {point}"
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*text_rgb)


# ---------- Main builder ----------

def build_pptx(plan: LecturePlan, theme: Dict[str, Any]) -> bytes:
    """
    Build a themed PPTX from LecturePlan.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    text_rgb = hex_to_rgb(theme["colors"]["text"])
    accent_rgb = hex_to_rgb(theme["colors"]["accent"])
    muted_rgb = hex_to_rgb(theme["colors"]["muted"])
    margins = theme["margins_in"]

    # Gradient background image for slides
    gradient_mode = "radial" if theme.get("header_bar") else "linear"
    bg_img = _render_gradient_bg(
        int(prs.slide_width.inches * 96),
        int(prs.slide_height.inches * 96),
        theme["background_gradient"],
        gradient_mode,
    )

    # ===== TITLE SLIDE =====
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_slide.shapes.add_picture(
        BytesIO(bg_img), 0, 0, prs.slide_width, prs.slide_height
    )

    # Top accent bar
    bar = title_slide.shapes.add_shape(
        1,  # rectangle
        0,
        0,
        prs.slide_width,
        Inches(0.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*accent_rgb)
    bar.line.fill.background()

    # Title
    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(2)
    )
    tf = title_box.text_frame
    tf.text = plan.topic
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(theme["sizes"]["title"] + 10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*text_rgb)
    try:
        p.font.name = theme["fonts"]["title"]
    except Exception:
        p.font.name = theme["fonts"].get("title_fallback", "Segoe UI")

    # Subtitle
    subtitle = f"{plan.language.upper()} • {plan.duration_minutes} minutes"
    sub_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(5.0), Inches(8), Inches(0.6)
    )
    sub_tf = sub_box.text_frame
    sub_tf.text = subtitle
    sp = sub_tf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(theme["sizes"]["subtitle"])
    sp.font.color.rgb = RGBColor(*muted_rgb)

    # Footer
    foot = title_slide.shapes.add_textbox(
        Inches(1), Inches(6.8), Inches(8), Inches(0.4)
    )
    foot_tf = foot.text_frame
    foot_tf.text = "EduSynth"
    fp = foot_tf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fp.font.size = Pt(theme["sizes"]["footer"])
    fp.font.color.rgb = RGBColor(*muted_rgb)

    # Bottom accent bar
    bar2 = title_slide.shapes.add_shape(
        1, 0, Inches(7.4), prs.slide_width, Inches(0.1)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = RGBColor(*accent_rgb)
    bar2.line.fill.background()

    # ===== CONTENT SLIDES =====
    for slide_item in plan.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            BytesIO(bg_img), 0, 0, prs.slide_width, prs.slide_height
        )

        # Corporate header bar if configured
        if theme.get("header_bar"):
            header_h = Inches(theme.get("header_bar_height_px", 90) / 96)
            header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, header_h)
            header.fill.solid()
            header.fill.fore_color.rgb = RGBColor(*accent_rgb)
            header.line.fill.background()

        # Title
        title_top = margins["top"] if not theme.get("header_bar") else 1.2
        tbox = slide.shapes.add_textbox(
            Inches(margins["left"]),
            Inches(title_top),
            Inches(10 - margins["left"] - margins["right"]),
            Inches(0.9),
        )
        ttf = tbox.text_frame
        ttf.text = slide_item.title
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.font.size = Pt(theme["sizes"]["title"])
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(*text_rgb)
        try:
            tp.font.name = theme["fonts"]["title"]
        except Exception:
            tp.font.name = theme["fonts"].get("title_fallback", "Segoe UI")

        # Accent rule under title
        rule = slide.shapes.add_shape(
            1,
            Inches(margins["left"]),
            Inches(title_top + 0.95),
            Inches(2),
            Inches(0.05),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(*accent_rgb)
        rule.line.fill.background()

        # Content area
        content_top = title_top + 1.2
        content_height = 7.5 - content_top - margins["bottom"]
        content_width = 10 - margins["left"] - margins["right"]

        layout_type = pick_layout(slide_item)
        theme_key = (plan.theme or "minimalist").lower()

        if layout_type == "process":
            add_process_diagram(
                slide.shapes,
                margins["left"],
                content_top,
                content_width,
                content_height - 0.2,
                slide_item.points,
                theme,
            )

        elif layout_type == "tree":
            add_tree_diagram(
                slide.shapes,
                margins["left"],
                content_top,
                content_width,
                content_height - 0.2,
                slide_item.points,
                theme,
            )

        elif layout_type == "timeline":
            add_timeline_diagram(
                slide.shapes,
                margins["left"],
                content_top + 0.5,
                content_width,
                content_height - 0.8,
                slide_item.points,
                theme,
            )

        elif layout_type == "compare":
            add_compare_layout(
                slide.shapes,
                margins["left"],
                content_top,
                content_width,
                content_height - 0.2,
                slide_item.points,
                theme,
            )

        else:
            # Text-only bullets
            get_point_icon(theme_key)
            tb = slide.shapes.add_textbox(
                Inches(margins["left"]),
                Inches(content_top),
                Inches(content_width),
                Inches(content_height),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            for i, point in enumerate(slide_item.points):
                para = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                para.text = f"{bullet_char} {point}"
                para.level = 0
                para.font.size = Pt(theme["sizes"]["body"])
                para.font.color.rgb = RGBColor(*text_rgb)
                try:
                    para.font.name = theme["fonts"]["body"]
                except Exception:
                    para.font.name = theme["fonts"].get("body_fallback", "Segoe UI")
                para.space_before = Pt(12)

        # Corporate footer: page number
        if theme.get("footer"):
            fb = slide.shapes.add_textbox(
                Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.3)
            )
            ft = fb.text_frame
            ft.text = str(slide_item.index + 1)
            fp = ft.paragraphs[0]
            fp.alignment = PP_ALIGN.RIGHT
            fp.font.size = Pt(theme["sizes"]["footer"])
            fp.font.color.rgb = RGBColor(*muted_rgb)

    # Save to bytes
    out = BytesIO()
    prs.save(out)
    return out.getvalue()
